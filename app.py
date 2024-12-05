import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import streamlit as st

# Load environment variables
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

# Configure Google Generative AI
genai.configure(api_key=GOOGLE_API_KEY)


# Function to extract text from PDF files
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


# Function to extract text from PPTX files
def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
    return text


# Function to split text into manageable chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)


# Function to create and save FAISS vector store
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")


# Function to load the conversational chain
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. 
    If the answer is not in the provided context, just say "answer is not available in the context".
    Do not provide a wrong answer.

    Context: {context}
    Question: {question}
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)


# Function to handle user input and respond to questions
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    # Safely load the FAISS vector store
    try:
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
    except ValueError as e:
        st.error(f"Error loading FAISS index: {e}")
        return

    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)

    st.write("Reply: ", response["output_text"])


# Main function to define the Streamlit app
def main():
    st.set_page_config(page_title="Chat with Documents", layout="wide")
    st.header("Chat with PDF and PPTX Files using Gemini 💁")

    user_question = st.text_input("Ask a Question from the Uploaded Documents")

    if user_question:
        user_input(user_question)

    with st.sidebar:
        st.title("Menu:")
        
        pdf_docs = st.file_uploader("Upload your PDF Files", accept_multiple_files=True, type=["pdf"])
        pptx_docs = st.file_uploader("Upload your PPTX Files", accept_multiple_files=True, type=["pptx"])

        if st.button("Submit & Process"):
            if pdf_docs or pptx_docs:
                with st.spinner("Processing..."):
                    raw_text = ""
                    if pdf_docs:
                        raw_text += get_pdf_text(pdf_docs)
                    if pptx_docs:
                        raw_text += get_pptx_text(pptx_docs)

                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("Documents have been processed and indexed.")
            else:
                st.error("Please upload at least one PDF or PPTX file.")


if __name__ == "__main__":
    main()
